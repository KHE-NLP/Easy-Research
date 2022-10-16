import io

from pdfminer.high_level import extract_text
import requests
from pptx import Presentation
from pptx.util import Inches
from PIL import Image


def get_pdf_text(pdf_data):
    if pdf_data[:8] == "https://" or pdf_data[:7] == "http://":
        resp = requests.get(pdf_data)
        resp_bytes = io.BytesIO(resp.content)
        data = extract_text(resp_bytes)
    else:
        data = open(pdf_data).read()
    return data


def get_pdf_paragraphs(pdf_data, min_par=40):
    data = get_pdf_text(pdf_data).split("\n\n")
    to_rem = []
    for i in data:
        if len(i) < min_par:
            to_rem.append(i)
    for i in to_rem:
        data.remove(i)
    return data


def get_pptx_slides(pptx_data, ask_desc=True):
    if pptx_data[:8] == "https://" or pptx_data[:7] == "http://":
        resp = requests.get(pptx_data)
        resp_bytes = io.BytesIO(resp.content)
        ppt = Presentation(resp_bytes)
    else:
        ppt = Presentation(pptx_data)

    slides = []
    for slide in ppt.slides:
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"
            if hasattr(shape, "image"):
                if ask_desc:
                    Image.open(io.BytesIO(shape.image.blob)).show()
                    desc = input("What is this? >")
                    slide_text += "[IMAGE " + desc + "]"
        slides.append(slide_text)
    return slides


def slice_well(string, place):
    for i in range(place, 0, -1):
        if string[i] == " " or string[i] == "\n":
            return [string[:i], string[i + 1:]]
    return [string]


def make_pptx_slides(pptx_strings):
    ppt = Presentation()
    layout = ppt.slide_layouts[0]

    for string in pptx_strings:
        slide = ppt.slides.add_slide(layout)
        title = string.split("\n")[0]
        lines = string.split("\n")[1:]
        new_lines = []
        for line in lines:
            tline = [line]
            # print(type(line))
            while len(tline[-1]) > 65:
                a = tline[-1]
                tline.remove(a)
                tline += slice_well(a, 65)
                # print(tline[-1])
            new_lines += tline
        text = "\n".join(new_lines)
        slide.shapes.title.text = title
        disp = Inches(1)
        sz = Inches(5)
        slide.shapes.title.top = disp
        body = slide.shapes.add_textbox(disp, 3 * disp, sz, sz)
        body.text_frame.text = text
        body.text_frame.left = disp
        body.text_frame.width = sz

    return ppt


if __name__ == "__main__":
    # pdf = get_pdf_paragraphs('http://faculty.cs.tamu.edu/stoleru/papers/won12re2mr.pdf')
    # print(*pdf, sep="\nPAR:\n")
    pptx = get_pptx_slides(
        "http://antares.cs.kent.edu/~seminar/Presentations/Minimum-Latency%20Broadcast%20Scheduling%20in%20Wireless%20Ad%20Hoc%20Networks.pptx",
        False)
    # print(*pptx, sep="\n")
    # pdata = ["Title\nMy data lol\nother line"]
    ppt = make_pptx_slides(pptx)
    ppt.save("test.pptx")
